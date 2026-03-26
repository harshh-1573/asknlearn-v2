import csv
import json
import os
from xml.etree import ElementTree

import docx
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation


TEXT_EXTENSIONS = {
    ".txt", ".md", ".rtf", ".csv", ".json", ".html", ".htm", ".xml",
    ".log", ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cpp",
    ".cs", ".go", ".php", ".sql", ".css", ".scss", ".less", ".yml", ".yaml",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif", ".gif"}


def _read_plain_text(path: str) -> str:
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
    for encoding in encodings:
        try:
            with open(path, "r", encoding=encoding) as file_handle:
                return file_handle.read()
        except UnicodeDecodeError:
            continue
    with open(path, "rb") as file_handle:
        return file_handle.read().decode("utf-8", errors="ignore")


def _read_csv(path: str) -> str:
    rows = []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as file_handle:
        reader = csv.reader(file_handle)
        for row in reader:
            if any(cell.strip() for cell in row):
                rows.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(rows)


def _read_json(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as file_handle:
        parsed = json.load(file_handle)
    return json.dumps(parsed, indent=2, ensure_ascii=False)


def _read_xml_like(path: str) -> str:
    tree = ElementTree.parse(path)
    text = " ".join(part.strip() for part in tree.itertext() if part and part.strip())
    if text:
        return text
    return _read_plain_text(path)


def _read_presentation(path: str) -> str:
    presentation = Presentation(path)
    chunks = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_chunks = [f"Slide {slide_index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                cleaned = shape.text.strip()
                if cleaned:
                    slide_chunks.append(cleaned)
        if len(slide_chunks) > 1:
            chunks.append("\n".join(slide_chunks))
    return "\n\n".join(chunks)


def _read_image(path: str) -> str:
    return pytesseract.image_to_string(Image.open(path))


def read_file(path: str, extension: str) -> str:
    extension = extension.lower()

    if extension == ".pdf":
        text = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text.append(page.extract_text() or "")
        return "\n".join(text).strip()

    if extension == ".docx":
        document = docx.Document(path)
        return "\n".join(para.text for para in document.paragraphs if para.text).strip()

    if extension in {".pptx", ".ppt"}:
        return _read_presentation(path).strip()

    if extension in IMAGE_EXTENSIONS:
        return _read_image(path).strip()

    if extension == ".csv":
        return _read_csv(path).strip()

    if extension == ".json":
        return _read_json(path).strip()

    if extension in {".xml", ".html", ".htm"}:
        return _read_xml_like(path).strip()

    if extension in TEXT_EXTENSIONS:
        return _read_plain_text(path).strip()

    return ""
